"""Export executive report to PDF, metrics to Excel, and a board deck to PowerPoint."""
import io
from datetime import datetime, timezone

import pandas as pd
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from app.models.models import FinancialMetric
from app.repositories.repositories import DashboardRepository


def export_executive_pdf(db) -> bytes:
    repo = DashboardRepository(db)
    executive = (repo.latest("executive") or type("x", (), {"payload": {}})).payload
    insights = (repo.latest("insights") or type("x", (), {"payload": {}})).payload

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=2 * cm)
    styles = getSampleStyleSheet()
    story = [
        Paragraph("Executive Financial Report", styles["Title"]),
        Paragraph(datetime.now(timezone.utc).strftime("Generated %d %b %Y %H:%M UTC"), styles["Normal"]),
        Spacer(1, 12),
        Paragraph(executive.get("headline", "Financial Overview"), styles["Heading2"]),
        Paragraph(str(executive.get("summary", "Run an analysis to populate this report.")), styles["BodyText"]),
        Spacer(1, 12),
    ]

    kpis = executive.get("kpis", {}) or {}
    rows = [["KPI", "Value"]] + [[k.replace("_", " ").title(), f"{v:,.0f}" if isinstance(v, (int, float)) else "—"]
                                 for k, v in kpis.items()]
    if len(rows) > 1:
        t = Table(rows, colWidths=[8 * cm, 6 * cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0f172a")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ]))
        story += [t, Spacer(1, 12)]

    for title, key in (("Key Highlights", "key_highlights"), ("Top Risks", "top_risks"),
                       ("Recommended Actions", "recommendations")):
        items = insights.get(key, [])
        if items:
            story.append(Paragraph(title, styles["Heading2"]))
            for item in items:
                text = item.get("action", str(item)) if isinstance(item, dict) else str(item)
                story.append(Paragraph(f"• {text}", styles["BodyText"]))
            story.append(Spacer(1, 8))

    doc.build(story)
    return buf.getvalue()


def export_metrics_excel(db) -> bytes:
    rows = db.query(FinancialMetric).all()
    df = pd.DataFrame([{"Fiscal Period": r.fiscal_period, "Metric": r.metric_name,
                        "Value": r.value, "Unit": r.unit, "Source": r.source} for r in rows])
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
        (df if not df.empty else pd.DataFrame(columns=["Fiscal Period", "Metric", "Value"])) \
            .to_excel(writer, index=False, sheet_name="Financial Metrics")
    return buf.getvalue()


_INK = RGBColor(0x0B, 0x12, 0x20)
_AMBER = RGBColor(0xFF, 0xB0, 0x20)
_MUTE = RGBColor(0x8C, 0xA0, 0xC6)
_WHITE = RGBColor(0xE9, 0xEC, 0xFA)


def _slide_bg(slide, prs) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _INK


def _add_title(slide, prs, text: str, size: int = 28) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), Inches(1))
    tf = box.text_frame
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = _WHITE


def _add_bullets(slide, prs, items: list[str], top: float = 1.5, color=_WHITE) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), prs.slide_width - Inches(1.2), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items or ["No data available — run an analysis first."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = color
        p.space_after = Pt(10)


def export_board_pptx(db) -> bytes:
    """Board-ready deck built from the same real dashboard data as the PDF/Excel
    exports — no fabricated slides, no content beyond what the analysis produced."""
    repo = DashboardRepository(db)
    executive = (repo.latest("executive") or type("x", (), {"payload": {}})).payload
    insights = (repo.latest("insights") or type("x", (), {"payload": {}})).payload

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — title
    s = prs.slides.add_slide(blank)
    _slide_bg(s, prs)
    _add_title(s, prs, "Executive Board Report", size=40)
    sub = s.shapes.add_textbox(Inches(0.6), Inches(1.5), prs.slide_width - Inches(1.2), Inches(1))
    sub.text_frame.text = executive.get("headline", "")
    sub.text_frame.paragraphs[0].font.size = Pt(20)
    sub.text_frame.paragraphs[0].font.color.rgb = _AMBER
    date_box = s.shapes.add_textbox(Inches(0.6), Inches(6.7), Inches(6), Inches(0.5))
    date_box.text_frame.text = datetime.now(timezone.utc).strftime("Generated %d %b %Y %H:%M UTC")
    date_box.text_frame.paragraphs[0].font.size = Pt(12)
    date_box.text_frame.paragraphs[0].font.color.rgb = _MUTE

    # Slide 2 — business health + KPIs
    s = prs.slides.add_slide(blank)
    _slide_bg(s, prs)
    _add_title(s, prs, f"Business Health Score: {executive.get('business_health_score', '—')}")
    kpis = executive.get("kpis", {}) or {}
    kpi_lines = [f"{k.replace('_', ' ').title()}: {v:,.2f}" if isinstance(v, (int, float)) else f"{k.replace('_', ' ').title()}: —"
                 for k, v in kpis.items()]
    _add_bullets(s, prs, kpi_lines)

    # Slide 3 — executive summary (conversational text, not bullets)
    s = prs.slides.add_slide(blank)
    _slide_bg(s, prs)
    _add_title(s, prs, "Executive Summary")
    body = s.shapes.add_textbox(Inches(0.6), Inches(1.5), prs.slide_width - Inches(1.2), Inches(5))
    body.text_frame.word_wrap = True
    body.text_frame.text = executive.get("summary") or "Run an analysis to populate this report."
    body.text_frame.paragraphs[0].font.size = Pt(16)
    body.text_frame.paragraphs[0].font.color.rgb = _WHITE

    # Slide 4 — green / red flags
    s = prs.slides.add_slide(blank)
    _slide_bg(s, prs)
    _add_title(s, prs, "Green Flags & Red Flags")
    half = prs.slide_width / 2
    green_box = s.shapes.add_textbox(Inches(0.6), Inches(1.5), half - Inches(0.9), Inches(5))
    green_box.text_frame.word_wrap = True
    for i, g in enumerate(insights.get("green_flags") or ["—"]):
        p = green_box.text_frame.paragraphs[0] if i == 0 else green_box.text_frame.add_paragraph()
        p.text = f"✓ {g}"; p.font.size = Pt(15); p.font.color.rgb = RGBColor(0x2D, 0xD4, 0xBF); p.space_after = Pt(10)
    red_box = s.shapes.add_textbox(Inches(0.6) + half, Inches(1.5), half - Inches(0.9), Inches(5))
    red_box.text_frame.word_wrap = True
    for i, r in enumerate(insights.get("red_flags") or ["—"]):
        p = red_box.text_frame.paragraphs[0] if i == 0 else red_box.text_frame.add_paragraph()
        p.text = f"✕ {r}"; p.font.size = Pt(15); p.font.color.rgb = RGBColor(0xF8, 0x71, 0x71); p.space_after = Pt(10)

    # Slide 5 — recommended actions
    s = prs.slides.add_slide(blank)
    _slide_bg(s, prs)
    _add_title(s, prs, "Recommended Actions")
    rec_lines = [f"{r.get('action', '')} — {r.get('priority', '')} priority, {r.get('timeframe', '')}"
                 for r in (insights.get("recommendations") or [])]
    _add_bullets(s, prs, rec_lines)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
